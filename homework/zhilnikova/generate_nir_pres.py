"""ПР5 по Жильниковой — Презентация к отчёту о НИР.
Корпоративный стиль ГУАП: синий #003F7F, красный акцент #C8102E,
шрифт Roboto, формат 16:9."""
import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Emu, Inches, Pt


OUT = os.path.join(os.path.dirname(__file__),
                   "Презентация_НИР_Рослов.pptx")

# ГУАП-палитра
GUAP_BLUE = RGBColor(0x00, 0x3F, 0x7F)
GUAP_RED = RGBColor(0xC8, 0x10, 0x2E)
TEXT_DARK = RGBColor(0x1A, 0x1A, 0x1A)
TEXT_GREY = RGBColor(0x55, 0x55, 0x55)
LIGHT_BG = RGBColor(0xF4, 0xF6, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Roboto"


def setup_slide_size(prs):
    # Стандартный размер PowerPoint 16:9 (12192000 × 6858000 EMU)
    prs.slide_width = 12192000
    prs.slide_height = 6858000
    # Убираем атрибут type="screen4x3" из <p:sldSz>, иначе LibreOffice
    # отказывается открывать файл (несоответствие размера и типа).
    sldSz = prs.part._element.find(
        "{http://schemas.openxmlformats.org/presentationml/2006/main}sldSz")
    if sldSz is not None and "type" in sldSz.attrib:
        del sldSz.attrib["type"]


def add_blank_slide(prs):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    return slide


def fill_solid(shape, color):
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color


def no_line(shape):
    shape.line.fill.background()


def add_rect(slide, x, y, w, h, fill_color, *, no_border=True):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    fill_solid(shp, fill_color)
    if no_border:
        no_line(shp)
    return shp


def add_textbox(slide, x, y, w, h, text, *,
                size=18, color=TEXT_DARK, bold=False, italic=False,
                align=PP_ALIGN.LEFT, font=FONT, line_spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, *,
                size=16, color=TEXT_DARK, font=FONT, bullet="•",
                space_after_pt=8, line_spacing=1.2):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after_pt)
        run = p.add_run()
        run.text = f"{bullet}  {item}"
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return tb


def add_header_band(slide, prs, title):
    """Тонкий синий бэнд сверху + красная полоса под ним."""
    add_rect(slide, 0, 0, prs.slide_width, Cm(2.6), GUAP_BLUE)
    add_rect(slide, 0, Cm(2.6), prs.slide_width, Cm(0.18), GUAP_RED)
    add_textbox(slide, Cm(1.2), Cm(0.55), prs.slide_width - Cm(2.4),
                Cm(1.8), title, size=24, color=WHITE, bold=True)


def add_footer(slide, prs, page_no, total):
    """Нижний колонтитул: страница + ГУАП | каф. №5."""
    y = prs.slide_height - Cm(0.9)
    add_rect(slide, 0, y, prs.slide_width, Cm(0.06), GUAP_BLUE)
    add_textbox(slide, Cm(1.2), y + Cm(0.15),
                Cm(20), Cm(0.6),
                "ГУАП • Кафедра № 5 • НИР, 2026",
                size=10, color=TEXT_GREY)
    add_textbox(slide,
                prs.slide_width - Cm(3.2), y + Cm(0.15),
                Cm(2.0), Cm(0.6),
                f"{page_no} / {total}", size=10, color=TEXT_GREY,
                align=PP_ALIGN.RIGHT)


# ============================================================
# Слайд 1 — Титульный
# ============================================================
def slide_title(prs):
    s = add_blank_slide(prs)
    # синий фон верхней половины
    add_rect(s, 0, 0, prs.slide_width, Cm(11.0), GUAP_BLUE)
    # красная полоса
    add_rect(s, 0, Cm(11.0), prs.slide_width, Cm(0.25), GUAP_RED)

    add_textbox(s, Cm(1.5), Cm(1.2), prs.slide_width - Cm(3.0),
                Cm(1.6),
                "САНКТ-ПЕТЕРБУРГСКИЙ ГОСУДАРСТВЕННЫЙ УНИВЕРСИТЕТ "
                "АЭРОКОСМИЧЕСКОГО ПРИБОРОСТРОЕНИЯ",
                size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(s, Cm(1.5), Cm(2.6), prs.slide_width - Cm(3.0),
                Cm(0.8),
                "Кафедра № 5 «Инноватики и интегрированных систем "
                "качества»", size=12, color=WHITE, italic=True,
                align=PP_ALIGN.CENTER)

    add_textbox(s, Cm(1.5), Cm(4.3), prs.slide_width - Cm(3.0),
                Cm(0.8),
                "ОТЧЁТ О НАУЧНО-ИССЛЕДОВАТЕЛЬСКОЙ РАБОТЕ",
                size=18, color=WHITE, bold=True,
                align=PP_ALIGN.CENTER)

    add_textbox(s, Cm(2.5), Cm(6.0), prs.slide_width - Cm(5.0),
                Cm(3.5),
                "Разработка системы поддержки принятия решений "
                "по оценке рисков внедрения технологий "
                "искусственного интеллекта в IT-компании на основе "
                "нечёткой логики",
                size=22, color=WHITE, bold=True,
                align=PP_ALIGN.CENTER, line_spacing=1.25)

    # Нижний блок — данные студента и руководителя
    add_textbox(s, Cm(1.8), Cm(12.4), Cm(15), Cm(0.7),
                "Магистрант:", size=12, color=GUAP_BLUE, bold=True)
    add_textbox(s, Cm(1.8), Cm(13.0), Cm(15), Cm(0.7),
                "Рослов Дмитрий Сергеевич", size=16, color=TEXT_DARK,
                bold=True)
    add_textbox(s, Cm(1.8), Cm(13.7), Cm(15), Cm(0.7),
                "группа М558М, направление 27.04.05 «Инноватика»",
                size=12, color=TEXT_GREY)

    add_textbox(s, Cm(18.2), Cm(12.4), Cm(14), Cm(0.7),
                "Научный руководитель:", size=12, color=GUAP_BLUE,
                bold=True)
    add_textbox(s, Cm(18.2), Cm(13.0), Cm(14), Cm(0.7),
                "Гетманова Галина Владимировна", size=16,
                color=TEXT_DARK, bold=True)
    add_textbox(s, Cm(18.2), Cm(13.7), Cm(14), Cm(0.7),
                "канд. экон. наук, доцент кафедры № 5",
                size=12, color=TEXT_GREY)

    add_textbox(s, Cm(1.8), Cm(15.4), Cm(15), Cm(0.7),
                "Дисциплина:", size=12, color=GUAP_BLUE, bold=True)
    add_textbox(s, Cm(1.8), Cm(16.0), Cm(15), Cm(0.7),
                "Научно-исследовательская работа", size=14,
                color=TEXT_DARK)

    add_textbox(s, Cm(18.2), Cm(15.4), Cm(14), Cm(0.7),
                "Преподаватель НИР:", size=12, color=GUAP_BLUE,
                bold=True)
    add_textbox(s, Cm(18.2), Cm(16.0), Cm(14), Cm(0.7),
                "Жильникова Н. А., профессор, д.т.н., доц.",
                size=14, color=TEXT_DARK)

    add_textbox(s, 0, prs.slide_height - Cm(1.5),
                prs.slide_width, Cm(0.8),
                "Санкт-Петербург, 2026", size=12, color=TEXT_GREY,
                align=PP_ALIGN.CENTER)


# ============================================================
# Слайд 2 — Актуальность и научно-практическая задача
# ============================================================
def slide_actuality(prs, page, total):
    s = add_blank_slide(prs)
    add_header_band(s, prs, "Актуальность темы и научно-практическая задача")
    add_footer(s, prs, page, total)

    # Левая колонка — актуальность
    add_textbox(s, Cm(1.2), Cm(3.5), Cm(15), Cm(0.8),
                "АКТУАЛЬНОСТЬ", size=14, color=GUAP_RED, bold=True)
    add_bullets(s, Cm(1.2), Cm(4.4), Cm(15.0), Cm(11),
        [
            "Рынок ИИ в РФ: 900 млрд ₽ в 2024 г. → 1,7 трлн ₽ к 2026 г.",
            "Указ Президента РФ № 490, ФЗ-123, Кодекс этики ИИ 2021 г.",
            "От 70 до 85 % ИИ-проектов не достигают целей (Gartner, IDC).",
            "Специфические риски: галлюцинации, concept drift, "
            "adversarial-атаки, регуляторные и этические ограничения.",
            "Классические FMEA / экспертные методы плохо работают "
            "при малых выборках и качественных оценках.",
        ], size=14, line_spacing=1.3, space_after_pt=10)

    # Правая колонка — научно-практическая задача
    add_rect(s, Cm(17.2), Cm(3.5), Cm(15.5), Cm(11.5), LIGHT_BG)
    add_textbox(s, Cm(17.7), Cm(3.8), Cm(15), Cm(0.8),
                "НАУЧНО-ПРАКТИЧЕСКАЯ ЗАДАЧА",
                size=14, color=GUAP_BLUE, bold=True)
    add_textbox(s, Cm(17.7), Cm(4.8), Cm(14.5), Cm(10),
                "Разработать систему поддержки принятия решений "
                "(СППР), обеспечивающую интегральную оценку рисков "
                "внедрения технологий ИИ в IT-компании путём "
                "интеграции модифицированного FMEA-анализа с "
                "нечётким выводом Мамдани и дефаззификацией методом "
                "центра тяжести, с категоризацией интегрального "
                "уровня риска R ∈ [0; 100] и формированием "
                "обоснованной рекомендации для лица, принимающего "
                "решение.",
                size=14, color=TEXT_DARK, line_spacing=1.35)


# ============================================================
# Слайд 3 — Цель и задачи
# ============================================================
def slide_goal(prs, page, total):
    s = add_blank_slide(prs)
    add_header_band(s, prs, "Цель и задачи магистерской диссертации")
    add_footer(s, prs, page, total)

    # Цель — выделенным блоком
    add_rect(s, Cm(1.2), Cm(3.5), prs.slide_width - Cm(2.4),
             Cm(2.6), LIGHT_BG)
    add_textbox(s, Cm(1.6), Cm(3.7), Cm(8), Cm(0.7),
                "ЦЕЛЬ", size=14, color=GUAP_RED, bold=True)
    add_textbox(s, Cm(1.6), Cm(4.4), prs.slide_width - Cm(3.2),
                Cm(2.0),
                "Разработать систему поддержки принятия решений по "
                "оценке рисков внедрения технологий искусственного "
                "интеллекта в IT-компании на основе аппарата "
                "нечёткой логики.",
                size=16, color=TEXT_DARK, line_spacing=1.3)

    # Задачи — нумерованный список
    add_textbox(s, Cm(1.2), Cm(6.6), Cm(20), Cm(0.8),
                "ЗАДАЧИ", size=14, color=GUAP_BLUE, bold=True)
    add_bullets(s, Cm(1.2), Cm(7.4), prs.slide_width - Cm(2.4),
                Cm(10),
        [
            "Проанализировать современное состояние ИИ-технологий в "
            "IT-секторе, идентифицировать и классифицировать "
            "характерные риски внедрения.",
            "Выполнить систематический обзор существующих методов "
            "оценки рисков и обосновать применимость аппарата "
            "нечёткой логики.",
            "Разработать гибридную модель оценки рисков: "
            "модифицированный FMEA + нечёткий вывод Мамдани + "
            "дефаззификация методом центра тяжести.",
            "Сформировать систему лингвистических переменных и базу "
            "продукционных правил, откалиброванную на "
            "ретроспективных данных IT-проектов.",
            "Реализовать СППР программно, провести верификацию и "
            "валидацию, выполнить технико-экономическое обоснование "
            "внедрения.",
        ], size=14, bullet="▸", space_after_pt=8)


# ============================================================
# Слайд 4 — Объект и предмет
# ============================================================
def slide_object(prs, page, total):
    s = add_blank_slide(prs)
    add_header_band(s, prs, "Объект и предмет исследования")
    add_footer(s, prs, page, total)

    # Два блока — объект и предмет
    block_w = Cm(15.4)
    block_h = Cm(10)
    y = Cm(4.5)

    # Объект
    add_rect(s, Cm(1.2), y, block_w, Cm(0.4), GUAP_BLUE)
    add_rect(s, Cm(1.2), y + Cm(0.4), block_w, block_h, LIGHT_BG)
    add_textbox(s, Cm(1.6), y + Cm(0.7), block_w - Cm(0.8), Cm(1),
                "ОБЪЕКТ ИССЛЕДОВАНИЯ",
                size=14, color=GUAP_BLUE, bold=True)
    add_textbox(s, Cm(1.6), y + Cm(2.0), block_w - Cm(0.8), Cm(8),
                "Процессы внедрения технологий искусственного "
                "интеллекта в IT-компаниях.",
                size=18, color=TEXT_DARK, line_spacing=1.3)

    # Предмет
    add_rect(s, Cm(17.3), y, block_w, Cm(0.4), GUAP_RED)
    add_rect(s, Cm(17.3), y + Cm(0.4), block_w, block_h, LIGHT_BG)
    add_textbox(s, Cm(17.7), y + Cm(0.7), block_w - Cm(0.8), Cm(1),
                "ПРЕДМЕТ ИССЛЕДОВАНИЯ",
                size=14, color=GUAP_RED, bold=True)
    add_textbox(s, Cm(17.7), y + Cm(2.0), block_w - Cm(0.8), Cm(8),
                "Методы и средства оценки рисков внедрения "
                "технологий искусственного интеллекта на основе "
                "аппарата нечёткой логики.",
                size=18, color=TEXT_DARK, line_spacing=1.3)


# ============================================================
# Слайд 5 — Выводы по главе 1
# ============================================================
def slide_conclusions(prs, page, total):
    s = add_blank_slide(prs)
    add_header_band(s, prs, "Выводы по первой главе")
    add_footer(s, prs, page, total)

    add_bullets(s, Cm(1.2), Cm(3.5), prs.slide_width - Cm(2.4),
                Cm(15),
        [
            "ИИ-технологии стали системообразующим фактором "
            "развития IT-сектора (российский рынок 650 → 1700 млрд ₽ "
            "к 2026 г., доля компаний с ИИ — 47 %).",
            "Внедрение ИИ сопряжено со специфическими рисками; "
            "предложена авторская типология из 6 категорий: "
            "технологические, безопасностные, организационно-кадровые, "
            "финансовые, регуляторно-этические, репутационные.",
            "Сформирована нормативная база (ГОСТ Р 31000-2019, "
            "ГОСТ Р 59276-2020, ISO/IEC 23894:2023, ISO/IEC 42001:2023, "
            "NIST AI RMF 1.0, EU AI Act 2024); типовые модели для "
            "IT-сектора отсутствуют.",
            "Классические методы (вероятностный анализ, дерево "
            "событий, классический FMEA) ограниченно применимы к "
            "ИИ-проектам с малыми выборками и качественной "
            "экспертной информацией.",
            "Аппарат нечёткой логики (Заде, Мамдани, Сугено) "
            "обеспечивает формализованную обработку качественной "
            "информации; гибридная схема «модифицированный FMEA + "
            "вывод Мамдани + центр тяжести» — методически "
            "обоснованный выбор.",
            "Сформулирована научно-практическая задача — "
            "разработать СППР с интегральной оценкой риска "
            "R ∈ [0; 100], категоризацией и обоснованной "
            "рекомендацией ЛПР.",
        ], size=13, bullet="✓", space_after_pt=10, line_spacing=1.25)


# ============================================================
# Слайд 6 — Апробация
# ============================================================
def slide_publications(prs, page, total):
    s = add_blank_slide(prs)
    add_header_band(s, prs, "Апробация результатов исследования")
    add_footer(s, prs, page, total)

    add_textbox(s, Cm(1.2), Cm(3.6), Cm(20), Cm(0.7),
                "ПУБЛИКАЦИИ ПО ТЕМЕ ИССЛЕДОВАНИЯ",
                size=14, color=GUAP_RED, bold=True)

    # Карточка статьи
    add_rect(s, Cm(1.2), Cm(4.6), prs.slide_width - Cm(2.4),
             Cm(5.6), LIGHT_BG)
    add_rect(s, Cm(1.2), Cm(4.6), Cm(0.18), Cm(5.6), GUAP_BLUE)
    add_textbox(s, Cm(1.7), Cm(4.85), Cm(28), Cm(0.7),
                "СТАТЬЯ № 1", size=12, color=GUAP_BLUE, bold=True)
    add_textbox(s, Cm(1.7), Cm(5.5), prs.slide_width - Cm(3.0),
                Cm(2.5),
                "Современное состояние и перспективы применения "
                "нечёткой логики для оценки рисков внедрения "
                "технологий искусственного интеллекта в "
                "IT-компаниях: систематический аналитический обзор",
                size=16, color=TEXT_DARK, bold=True, line_spacing=1.25)
    add_textbox(s, Cm(1.7), Cm(8.4), prs.slide_width - Cm(3.0),
                Cm(1.4),
                "Рослов Д. С. — Санкт-Петербург : ГУАП, 2026. — "
                "16 с. (НИР магистра, кафедра № 5)\n"
                "Объём: 22 927 знаков, 24 источника, 3 рисунка, "
                "3 таблицы, PRISMA-обзор",
                size=12, color=TEXT_GREY, line_spacing=1.4)

    # Конференции
    add_textbox(s, Cm(1.2), Cm(11.0), Cm(20), Cm(0.7),
                "ПЛАНИРУЕМЫЕ ВЫСТУПЛЕНИЯ И ПУБЛИКАЦИИ",
                size=14, color=GUAP_BLUE, bold=True)
    add_bullets(s, Cm(1.2), Cm(11.9), prs.slide_width - Cm(2.4),
                Cm(6),
        [
            "Ежегодная научно-техническая конференция ГУАП — "
            "Санкт-Петербург, 2027 г.",
            "Международная научно-практическая конференция "
            "«Инновационная экономика и менеджмент: методы и "
            "технологии» — Санкт-Петербург, 2027 г.",
            "Публикации в изданиях из перечня ВАК и базы РИНЦ — "
            "2026–2027 уч. г.",
        ], size=14, bullet="•", space_after_pt=8)


# ============================================================
# Слайд 7 — Заключение
# ============================================================
def slide_conclusion(prs, page, total):
    s = add_blank_slide(prs)
    add_header_band(s, prs, "Заключение и дальнейшие шаги")
    add_footer(s, prs, page, total)

    # Краткие выводы
    add_textbox(s, Cm(1.2), Cm(3.5), Cm(20), Cm(0.7),
                "ОСНОВНЫЕ РЕЗУЛЬТАТЫ НИР",
                size=14, color=GUAP_RED, bold=True)
    add_bullets(s, Cm(1.2), Cm(4.3), Cm(15.4), Cm(11),
        [
            "Сформирован теоретический базис исследования.",
            "Предложена типология рисков ИИ из 6 категорий.",
            "Проанализирована нормативно-правовая база РФ и мира.",
            "Обосновано применение аппарата нечёткой логики и "
            "схемы «FMEA + Мамдани + центр тяжести».",
            "Сформулирована научно-практическая задача.",
            "Подготовлена научная статья по обзору (16 с., "
            "24 источника).",
        ], size=14, bullet="✓", space_after_pt=8, line_spacing=1.25)

    # План на следующие этапы
    add_rect(s, Cm(17.3), Cm(3.5), Cm(15.4), Cm(11.7), LIGHT_BG)
    add_textbox(s, Cm(17.7), Cm(3.8), Cm(15), Cm(0.7),
                "ДАЛЬНЕЙШИЕ ШАГИ ПО ВКРМ",
                size=14, color=GUAP_BLUE, bold=True)
    add_bullets(s, Cm(17.7), Cm(4.6), Cm(14.6), Cm(11),
        [
            "Сент.–дек. 2026 г. — глава 2: ретроспективный анализ "
            "ИИ-проектов в IT-компаниях, выявление пробелов.",
            "Янв.–апр. 2027 г. — глава 3: разработка архитектуры "
            "СППР, лингвистических переменных, базы правил.",
            "Апр.–май 2027 г. — программная реализация СППР, "
            "верификация и валидация на ретроспективных данных, "
            "ТЭО.",
            "Май 2027 г. — оформление по ГОСТ, нормоконтроль, "
            "проверка на оригинальность.",
            "Июнь 2027 г. — предзащита, защита ВКРМ.",
        ], size=13, bullet="▸", space_after_pt=8, line_spacing=1.25)


def main():
    prs = Presentation()
    setup_slide_size(prs)

    total = 7
    slide_title(prs)
    slide_actuality(prs, 2, total)
    slide_goal(prs, 3, total)
    slide_object(prs, 4, total)
    slide_conclusions(prs, 5, total)
    slide_publications(prs, 6, total)
    slide_conclusion(prs, 7, total)

    prs.save(OUT)
    print(f"OK: {OUT}")


if __name__ == "__main__":
    main()
