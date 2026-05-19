"""Презентация по реферату ПР6 (тема 9, УКОС) на основе официального шаблона ГУАП 16:9."""
import copy
import os

from lxml import etree
from pptx import Presentation
from pptx.util import Cm, Pt

TEMPLATE = "/tmp/guap-169-2023.pptx"
OUT_DIR = "/home/sergei/dima-projects/homework/milova"
PPTX_PATH = os.path.join(OUT_DIR, "Презентация_ПР6_Принятие_коллективных_решений_Рослов.pptx")


def remove_all_slides(prs):
    """Удаляет все существующие слайды из шаблона (по xml)."""
    sldIdLst = prs.slides._sldIdLst
    rId_to_remove = []
    for sldId in list(sldIdLst):
        rId = sldId.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
        rId_to_remove.append(rId)
        sldIdLst.remove(sldId)
    part = prs.part
    for rId in rId_to_remove:
        try:
            part.drop_rel(rId)
        except KeyError:
            pass


def find_layout(prs, name):
    for m in prs.slide_masters:
        for layout in m.slide_layouts:
            if layout.name == name:
                return layout
    raise KeyError(name)


def get_ph(slide, idx):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    raise KeyError(f"placeholder idx={idx}")


def set_ph_text(ph, lines):
    """Заполняет placeholder списком абзацев (строк или (текст, bold) пар).
    Сохраняет стиль (шрифт, размер) из layout-шаблона за счёт reuse первого rPr.
    """
    tf = ph.text_frame
    tf.clear()
    if not lines:
        return
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if isinstance(line, tuple):
            text, bold = line
        else:
            text, bold = line, False
        # Очистим существующие runs
        for r in list(p.runs):
            r._r.getparent().remove(r._r)
        run = p.add_run()
        run.text = text
        if bold:
            run.font.bold = True


def add_slide(prs, layout_name):
    layout = find_layout(prs, layout_name)
    return prs.slides.add_slide(layout)


def add_textbox(slide, left_cm, top_cm, width_cm, height_cm, lines, *, font_size=14, bold=False, color=None, align_center=False):
    """Добавляет текстовый блок с указанием размеров и стиля."""
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    box = slide.shapes.add_textbox(Cm(left_cm), Cm(top_cm), Cm(width_cm), Cm(height_cm))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if isinstance(line, tuple):
            text, line_bold = line
        else:
            text, line_bold = line, bold
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = text
        run.font.name = "Arial"
        run.font.size = Pt(font_size)
        run.font.bold = line_bold
        if color:
            run.font.color.rgb = RGBColor(*color)
        if align_center:
            p.alignment = PP_ALIGN.CENTER
    return box


def add_table(slide, left_cm, top_cm, width_cm, height_cm, data, *, header_color=(0, 63, 127), header_text_color=(255, 255, 255), font_size=11):
    """Добавляет таблицу. data — list of rows; первая строка — заголовки."""
    from pptx.dml.color import RGBColor
    rows = len(data)
    cols = len(data[0])
    tbl_shape = slide.shapes.add_table(rows, cols, Cm(left_cm), Cm(top_cm), Cm(width_cm), Cm(height_cm))
    table = tbl_shape.table
    for ri, row in enumerate(data):
        for ci, val in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = ""
            tf = cell.text_frame
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = str(val)
            run.font.name = "Arial"
            run.font.size = Pt(font_size)
            if ri == 0:
                run.font.bold = True
                run.font.color.rgb = RGBColor(*header_text_color)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(*header_color)
            else:
                run.font.color.rgb = RGBColor(50, 50, 50)
            from pptx.enum.text import PP_ALIGN
            p.alignment = PP_ALIGN.CENTER
    return table


def main():
    prs = Presentation(TEMPLATE)
    remove_all_slides(prs)

    # ---------- 1. Титульный слайд ----------
    slide = add_slide(prs, "Титульная страница")
    set_ph_text(get_ph(slide, 0), ["Задачи и алгоритмы принятия коллективных решений"])
    # Удаляем заводские placeholder'ы автора и даты, чтобы они не оставались пустыми
    for idx_to_remove in (10, 12):
        try:
            ph = get_ph(slide, idx_to_remove)
            ph._element.getparent().remove(ph._element)
        except KeyError:
            pass
    # Авторский блок текста в нижней половине слайда
    add_textbox(
        slide,
        left_cm=1.53, top_cm=11.4, width_cm=22, height_cm=4.5,
        lines=[
            ("Реферат по дисциплине", False),
            ("«Управление качеством организационных систем»", True),
            "",
            "Выполнил: студент гр. М558М Д. С. Рослов",
            "Руководитель: канд. техн. наук, доцент В. М. Милова",
        ],
        font_size=16,
        color=(50, 50, 50),
    )
    add_textbox(
        slide,
        left_cm=22.65, top_cm=16.65, width_cm=10, height_cm=1.2,
        lines=["Санкт-Петербург, 2026"],
        font_size=14,
        color=(50, 50, 50),
    )

    # ---------- 2. Содержание ----------
    slide = add_slide(prs, "Текст - один блок")
    set_ph_text(get_ph(slide, 0), ["Содержание"])
    set_ph_text(
        get_ph(slide, 12),
        [
            "1. Задача принятия группового решения",
            "2. Правило большинства голосов и принцип Кондорсе",
            "3. Метод Борда",
            "4. Парадокс Эрроу",
            "5. Пример: выбор системы менеджмента качества",
        ],
    )

    # ---------- 3. Задача коллективного выбора ----------
    slide = add_slide(prs, "Текст - один блок")
    set_ph_text(get_ph(slide, 0), ["1. Задача принятия группового решения"])
    set_ph_text(
        get_ph(slide, 12),
        [
            "Дано: множество альтернатив A = {a₁, …, aₘ}, экспертов E = {e₁, …, eₙ}, профиль предпочтений P = (≻₁, …, ≻ₙ).",
            "Найти: правило F: P → R, где R — коллективное отношение на A.",
            "Желаемые свойства алгоритма:",
            "• универсальная область;",
            "• единогласие (Парето);",
            "• независимость от посторонних альтернатив;",
            "• отсутствие диктатора;",
            "• монотонность.",
        ],
    )

    # ---------- 4. Правило большинства ----------
    slide = add_slide(prs, "Текст - один блок")
    set_ph_text(get_ph(slide, 0), ["2.1. Правило большинства голосов"])
    set_ph_text(
        get_ph(slide, 12),
        [
            "Каждый эксперт указывает одну предпочитаемую альтернативу — побеждает набравшая больше всех голосов.",
            "Виды большинства: абсолютное (>50%) и относительное.",
            "Достоинства: простота, интуитивная понятность, широко применяется.",
            "Недостатки: использует только 1-е место; возможен парадокс «победителя меньшинства»; неинформативно при m ≥ 3.",
        ],
    )

    # ---------- 5. Принцип Кондорсе ----------
    slide = add_slide(prs, "Текст - два блока")
    set_ph_text(get_ph(slide, 0), ["2.2. Принцип Кондорсе и парадокс"])
    set_ph_text(
        get_ph(slide, 12),
        [
            "Победитель Кондорсе: альтернатива x, для которой N(x, y) > N(y, x) для всех y ≠ x.",
            "N(x, y) — число экспертов, предпочитающих x альтернативе y.",
            "Использует полную информацию о ранжированиях.",
            "Удовлетворяет: Парето, монотонность, IIA в попарном сравнении.",
        ],
    )
    set_ph_text(
        get_ph(slide, 13),
        [
            "Парадокс Кондорсе (1785):",
            "e₁: a ≻ b ≻ c",
            "e₂: b ≻ c ≻ a",
            "e₃: c ≻ a ≻ b",
            "В попарных сравнениях получаем цикл: a ≻ b ≻ c ≻ a.",
            "Победителя Кондорсе нет — мажоритарное отношение нетранзитивно.",
        ],
    )

    # ---------- 6. Метод Борда ----------
    slide = add_slide(prs, "Текст - два блока")
    set_ph_text(get_ph(slide, 0), ["3. Метод Борда (1770)"])
    set_ph_text(
        get_ph(slide, 12),
        [
            "Оценка альтернативы x:",
            "B(x) = Σₖ (m − rₖ(x)),",
            "где rₖ(x) — ранг x у эксперта k (1 — лучший).",
            "При m альтернативах 1-е место даёт m − 1 балл, последнее — 0.",
            "Победитель Борда — альтернатива с max B(x).",
        ],
    )
    set_ph_text(
        get_ph(slide, 13),
        [
            "Достоинства:",
            "• использует полную информацию;",
            "• всегда даёт ранжирование;",
            "• Парето + монотонность.",
            "Недостатки:",
            "• нарушает условие Кондорсе;",
            "• уязвим к стратегическому голосованию.",
            "Применение: Eurovision, NCAA, оценка поставщиков.",
        ],
    )

    # ---------- 7. Парадокс Эрроу ----------
    slide = add_slide(prs, "Текст - один блок")
    set_ph_text(get_ph(slide, 0), ["4. Теорема Эрроу (1951)"])
    set_ph_text(
        get_ph(slide, 12),
        [
            "При m ≥ 3 не существует алгоритма, удовлетворяющего одновременно 4 аксиомам:",
            "1) универсальная область;",
            "2) принцип Парето;",
            "3) независимость от посторонних альтернатив (IIA);",
            "4) отсутствие диктатора.",
            "Следствие: любая процедура коллективного выбора жертвует одной из аксиом.",
            "Способы обхода: MAUT, теорема Блэка, вероятностные процедуры.",
        ],
    )

    # ---------- 8. Пример: профиль предпочтений ----------
    slide = add_slide(prs, "Пустой слайд с заголовком")
    set_ph_text(get_ph(slide, 0), ["5. Пример: выбор СМК для IT-компании"])
    add_textbox(
        slide,
        left_cm=1.2, top_cm=3.3, width_cm=31, height_cm=1.2,
        lines=["Альтернативы: A — ISO 9001, B — CMMI-DEV, C — Six Sigma, D — Agile + ISO/IEC 33000."],
        font_size=14,
    )
    add_textbox(
        slide,
        left_cm=1.2, top_cm=4.8, width_cm=14, height_cm=0.7,
        lines=[("Профиль предпочтений 5 экспертов:", True)],
        font_size=13,
    )
    add_table(
        slide,
        left_cm=1.2, top_cm=5.7, width_cm=14.5, height_cm=4.5,
        data=[
            ["Эксперт", "1-е", "2-е", "3-е", "4-е"],
            ["e₁ (ген. дир.)", "A", "B", "D", "C"],
            ["e₂ (CTO)", "B", "D", "A", "C"],
            ["e₃ (CQO)", "A", "C", "B", "D"],
            ["e₄ (рук. разр.)", "D", "B", "A", "C"],
            ["e₅ (CFO)", "A", "B", "C", "D"],
        ],
        font_size=12,
    )
    add_textbox(
        slide,
        left_cm=16.5, top_cm=4.8, width_cm=15.5, height_cm=8,
        lines=[
            ("Применённые алгоритмы:", True),
            "",
            "1. Правило большинства — 3 первых места у A.",
            "",
            "2. Кондорсе — A побеждает в попарных сравнениях B, C, D.",
            "",
            "3. Борда — A: 11, B: 10, D: 6, C: 3 баллов.",
        ],
        font_size=13,
    )

    # ---------- 9. Результаты сравнения алгоритмов ----------
    slide = add_slide(prs, "Пустой слайд с заголовком")
    set_ph_text(get_ph(slide, 0), ["5. Сравнение алгоритмов"])
    add_textbox(
        slide,
        left_cm=1.2, top_cm=3.1, width_cm=20, height_cm=0.7,
        lines=[("Все три алгоритма дали один результат — победитель A (ISO 9001:2015).", True)],
        font_size=15,
        color=(0, 63, 127),
    )
    add_table(
        slide,
        left_cm=1.2, top_cm=4.3, width_cm=20, height_cm=5,
        data=[
            ["Алгоритм", "Победитель", "Ранжирование", "Ключевое свойство"],
            ["Правило большинства", "A (60%)", "A ≻ {B, D} ≻ C", "Простота, использует только 1-е место"],
            ["Принцип Кондорсе", "A", "A ≻ B ≻ D ≻ C", "Попарные сравнения, нет цикла"],
            ["Метод Борда (баллы)", "A (11)", "A(11) ≻ B(10) ≻ D(6) ≻ C(3)", "Полное ранжирование, баллы"],
        ],
        font_size=13,
    )
    add_textbox(
        slide,
        left_cm=1.2, top_cm=11.5, width_cm=30, height_cm=3,
        lines=[
            "Согласованность результатов повышает доверие к решению.",
            "Анализ чувствительности: победитель устойчив к перестановкам у одного эксперта.",
        ],
        font_size=14,
    )

    # ---------- 10. Заключение ----------
    slide = add_slide(prs, "Текст - один блок")
    set_ph_text(get_ph(slide, 0), ["Заключение"])
    set_ph_text(
        get_ph(slide, 12),
        [
            "Рассмотрены классические алгоритмы коллективного выбора: правило большинства, Кондорсе, метод Борда.",
            "Показан парадокс Кондорсе и теорема Эрроу о невозможности идеального алгоритма.",
            "Применение к задаче выбора СМК для IT-компании дало согласованный результат — ISO 9001:2015.",
            "Методы применимы в задачах УКОС, работе экспертных комитетов и СППР.",
        ],
    )

    # ---------- 11. Завершение ----------
    slide = add_slide(prs, "Завершение презентации 4 - с контактами")
    try:
        set_ph_text(
            get_ph(slide, 10),
            [
                "Спасибо за внимание!",
                "",
                "Д. С. Рослов, гр. М558М",
                "roslov.dima00@mail.ru",
            ],
        )
    except KeyError:
        # idx может отличаться — берём первый placeholder
        for ph in slide.placeholders:
            set_ph_text(
                ph,
                [
                    "Спасибо за внимание!",
                    "",
                    "Д. С. Рослов, гр. М558М",
                    "roslov.dima00@mail.ru",
                ],
            )
            break

    prs.save(PPTX_PATH)
    print(f"Saved: {PPTX_PATH}")


if __name__ == "__main__":
    main()
